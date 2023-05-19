from pptx import Presentation
from PIL import Image
import os

def ppt_to_txt(ppt_path, output_dir):
    # Open the ppt file
    prs = Presentation(ppt_path)

    # Get the filename of the ppt (excluding the extension)
    base_name = os.path.basename(ppt_path)
    name_without_ext = os.path.splitext(base_name)[0]

    # Create the full path for the txt file
    txt_path = os.path.join(output_dir, f'{name_without_ext}_2txt.txt')

    # Write the text content from the ppt to the txt file
    with open(txt_path, 'w', encoding='utf-8') as f:
        # Iterate through each slide
        for slide in prs.slides:
            # Iterate through each shape in a slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Iterate through each paragraph in a text frame
                    for paragraph in shape.text_frame.paragraphs:
                        # Iterate through each run (sentence) in a paragraph
                        for run in paragraph.runs:
                            # Write the sentence
                            f.write(run.text)
                        # Add a new line after each paragraph
                        f.write('\n')
